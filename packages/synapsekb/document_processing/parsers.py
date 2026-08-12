from __future__ import annotations

import html
from dataclasses import dataclass
from pathlib import Path
from zipfile import BadZipFile

import fitz
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation


class NeedsOcrError(RuntimeError):
    pass


class UnsupportedDocumentError(ValueError):
    pass


@dataclass(frozen=True, slots=True)
class ParsedPage:
    number: int
    markdown: str


@dataclass(frozen=True, slots=True)
class ParsedDocument:
    markdown: str
    pages: list[ParsedPage]

    @property
    def page_count(self) -> int:
        return len(self.pages)


def _parse_pdf(path: Path) -> ParsedDocument:
    pdf = fitz.open(path)
    pages: list[ParsedPage] = []
    total_chars = 0
    for index, page in enumerate(pdf):
        text = page.get_text("text").strip()
        total_chars += len(text)
        pages.append(ParsedPage(index + 1, text))
    if pages and total_chars / len(pages) < 40:
        raise NeedsOcrError("PDF 每页可提取文字过少，需要 OCR")
    markdown = "\n\n".join(f"<!-- page:{page.number} -->\n{page.markdown}" for page in pages)
    return ParsedDocument(markdown=markdown, pages=pages)


def _parse_docx(path: Path) -> ParsedDocument:
    document = DocxDocument(str(path))
    lines: list[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style = paragraph.style.name.lower() if paragraph.style else ""
        if style.startswith("heading"):
            level = next((part for part in style.split() if part.isdigit()), "2")
            lines.append(f"{'#' * min(int(level), 6)} {text}")
        else:
            lines.append(text)
    markdown = "\n\n".join(lines)
    return ParsedDocument(markdown=markdown, pages=[ParsedPage(1, markdown)])


def _parse_xlsx(path: Path) -> ParsedDocument:
    workbook = load_workbook(path, read_only=True, data_only=True)
    sections: list[str] = []
    for worksheet in workbook.worksheets:
        lines = [f"## {worksheet.title}"]
        for row in worksheet.iter_rows(values_only=True):
            cells = ["" if cell is None else str(cell) for cell in row]
            if any(cells):
                lines.append(" | ".join(cells))
        sections.append("\n".join(lines))
    markdown = "\n\n".join(sections)
    return ParsedDocument(markdown=markdown, pages=[ParsedPage(1, markdown)])


def _parse_pptx(path: Path) -> ParsedDocument:
    presentation = Presentation(str(path))
    pages: list[ParsedPage] = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        pages.append(ParsedPage(index, "\n\n".join(texts)))
    markdown = "\n\n".join(f"<!-- page:{p.number} -->\n{p.markdown}" for p in pages)
    return ParsedDocument(markdown=markdown, pages=pages)


def _parse_html(path: Path) -> ParsedDocument:
    raw = path.read_text(encoding="utf-8", errors="replace")
    soup = BeautifulSoup(raw, "html.parser")
    for element in soup(["script", "style", "noscript"]):
        element.decompose()
    text = "\n".join(line.strip() for line in soup.get_text("\n").splitlines() if line.strip())
    return ParsedDocument(markdown=html.unescape(text), pages=[ParsedPage(1, text)])


def parse_document(path: Path, filename: str, media_type: str) -> ParsedDocument:
    suffix = Path(filename).suffix.lower()
    try:
        if suffix in {".jpg", ".jpeg", ".png", ".tif", ".tiff"} or media_type.startswith("image/"):
            raise NeedsOcrError("图片文档需要 OCR")
        if suffix == ".pdf" or media_type == "application/pdf":
            return _parse_pdf(path)
        if suffix == ".docx":
            return _parse_docx(path)
        if suffix == ".xlsx":
            return _parse_xlsx(path)
        if suffix == ".pptx":
            return _parse_pptx(path)
        if suffix in {".html", ".htm"} or media_type == "text/html":
            return _parse_html(path)
        if suffix in {".md", ".markdown", ".txt"} or media_type.startswith("text/"):
            text = path.read_text(encoding="utf-8", errors="replace")
            return ParsedDocument(markdown=text, pages=[ParsedPage(1, text)])
    except (BadZipFile, ValueError) as exc:
        raise UnsupportedDocumentError(f"文档结构损坏: {filename}") from exc
    raise UnsupportedDocumentError(f"不支持的文件类型: {suffix or media_type}")
